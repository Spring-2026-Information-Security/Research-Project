"""Build Presentation.pptx - a 15-minute deck mirroring PRESENTATION.md.

Produces a clean, basic-themed deck with:
* Title slide and section-divider slides
* Bullet content slides
* Tables rendered as native PowerPoint tables
* Monospace code blocks
* Embedded chart PNGs from the latest benchmark run

Run from the repo root:

    python scripts/build_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = REPO_ROOT / "Presentation.pptx"

# ── theme ───────────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x3B, 0x5C)
NAVY_LIGHT = RGBColor(0x2F, 0x5C, 0x8A)
ACCENT_RED = RGBColor(0xD3, 0x36, 0x36)
ACCENT_GREEN = RGBColor(0x2D, 0x8F, 0x5A)
ACCENT_AMBER = RGBColor(0xC8, 0x7A, 0x00)
TEXT_DARK = RGBColor(0x1F, 0x1F, 0x1F)
TEXT_MUTED = RGBColor(0x6A, 0x6A, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
TABLE_HEADER_BG = NAVY
ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)

FONT_HEADING = "Calibri"
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

CHART_DIR = REPO_ROOT / "login-lab" / "logs" / "benchmark"

# Pin to the 100-word, 50-trial run (the only complete run across all 21 configs).
# The 200/300/500-word runs are referenced via the cross-run scaling charts that
# were also written into this dir by scripts/chart_wordlist_scaling.py.
HEADLINE_RUN = "20260428T164449Z"


# ── helpers ────────────────────────────────────────────────────────────────


def headline_run_dir() -> Path | None:
    p = CHART_DIR / HEADLINE_RUN
    return p if p.exists() else None


LATEST = headline_run_dir()


def chart(name: str) -> Path | None:
    if LATEST is None:
        return None
    p = LATEST / name
    return p if p.exists() else None


def add_blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def fill_solid(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    fill_solid(shape, color)
    shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, *, font_size=18, bold=False,
                color=TEXT_DARK, font=FONT_BODY, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title_bar(slide, slide_width: Emu, title: str, subtitle: str | None = None) -> None:
    add_rect(slide, 0, 0, slide_width, Inches(0.10), NAVY)
    add_textbox(
        slide, Inches(0.5), Inches(0.25), slide_width - Inches(1.0), Inches(0.7),
        title, font_size=30, bold=True, color=NAVY, font=FONT_HEADING,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.5), Inches(0.95), slide_width - Inches(1.0), Inches(0.4),
            subtitle, font_size=16, color=TEXT_MUTED, font=FONT_BODY,
        )


def add_footer(slide, slide_width: Emu, slide_height: Emu, page: int, total: int, label: str) -> None:
    add_rect(slide, 0, slide_height - Inches(0.08), slide_width, Inches(0.08), NAVY)
    add_textbox(
        slide, Inches(0.4), slide_height - Inches(0.42),
        slide_width - Inches(0.8), Inches(0.3),
        label, font_size=10, color=TEXT_MUTED,
    )
    add_textbox(
        slide, slide_width - Inches(1.5), slide_height - Inches(0.42),
        Inches(1.1), Inches(0.3),
        f"{page} / {total}", font_size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
    )


def add_bullets(slide, x, y, w, h, items, *, font_size=18, indent_per_level=0.25,
                bullet_color=NAVY, text_color=TEXT_DARK, line_spacing=1.15):
    """items: list of (level, text) or list of plain strings (level=0)."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.05)

    first = True
    for entry in items:
        if isinstance(entry, tuple):
            level, text = entry
        else:
            level, text = 0, entry
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.line_spacing = line_spacing
        p.alignment = PP_ALIGN.LEFT
        # Bullet glyph
        bullet = "• " if level == 0 else "- "
        run_b = p.add_run()
        run_b.text = bullet
        run_b.font.name = FONT_BODY
        run_b.font.size = Pt(font_size)
        run_b.font.bold = True
        run_b.font.color.rgb = bullet_color
        # Text
        run_t = p.add_run()
        run_t.text = text
        run_t.font.name = FONT_BODY
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = text_color
    return box


def add_code_block(slide, x, y, w, h, code: str, *, font_size=14):
    box = slide.shapes.add_textbox(x, y, w, h)
    fill_solid(box, LIGHT_BG)
    box.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    lines = code.rstrip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_DARK
    return box


def add_table(slide, x, y, w, h, headers, rows, *, font_size=12, header_color=TABLE_HEADER_BG):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = tbl_shape.table
    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        fill_solid(cell, header_color)
        cell.text = ""
        tf = cell.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header
        run.font.name = FONT_BODY
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = WHITE
    # Body rows
    for i, row in enumerate(rows, start=1):
        bg = ROW_ALT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            fill_solid(cell, bg)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            # Support tuples (text, color) for accent coloring
            if isinstance(val, tuple):
                text, color = val
            else:
                text, color = val, TEXT_DARK
            run = p.add_run()
            run.text = str(text)
            run.font.name = FONT_BODY
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return tbl_shape


def add_image_safe(slide, png_path: Path | None, x, y, w, h, fallback_caption: str = "(chart not generated)"):
    if png_path and png_path.exists():
        slide.shapes.add_picture(str(png_path), x, y, width=w, height=h)
    else:
        add_rect(slide, x, y, w, h, LIGHT_BG)
        add_textbox(slide, x, y, w, h, fallback_caption,
                    font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── slide builders ──────────────────────────────────────────────────────────


def slide_title(prs, w, h):
    s = add_blank_slide(prs)
    # Top color band
    add_rect(s, 0, 0, w, Inches(2.2), NAVY)
    add_rect(s, 0, Inches(2.2), w, Inches(0.06), ACCENT_RED)
    # Title block
    add_textbox(s, Inches(0.6), Inches(0.7), w - Inches(1.2), Inches(0.7),
                "CS 47205/57205 · Project 3", font_size=18, color=WHITE, font=FONT_BODY)
    add_textbox(s, Inches(0.6), Inches(1.05), w - Inches(1.2), Inches(1.0),
                "Measuring Online Password Guessing Resistance",
                font_size=36, bold=True, color=WHITE, font=FONT_HEADING)
    # Subtitle below band
    add_textbox(s, Inches(0.6), Inches(2.6), w - Inches(1.2), Inches(0.6),
                "A reproducible measurement framework for authentication defenses",
                font_size=22, color=NAVY, font=FONT_HEADING)
    add_textbox(s, Inches(0.6), Inches(3.4), w - Inches(1.2), Inches(0.5),
                "Austin & Ian", font_size=18, color=TEXT_DARK, font=FONT_BODY)


def slide_problem(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "The problem", "Why measure password-guessing defenses systematically")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(2.5), [
        "Online password guessing is the single most common identity attack on the public Internet.",
        "Every authentication system ships a different mix of defenses:",
        (1, "account lockout, IP rate limiting, progressive delays, CAPTCHAs, MFA, anomaly detection, bot filters …"),
        "These defenses interact in non-obvious ways and are often misconfigured.",
    ], font_size=18)
    add_rect(s, Inches(0.6), Inches(4.2), w - Inches(1.2), Inches(0.06), ACCENT_RED)
    add_bullets(s, Inches(0.6), Inches(4.4), w - Inches(1.2), Inches(2.0), [
        ("Question: how do these defenses actually compare under controlled, repeatable measurement?"),
        ("Goal: build attacker + target framework, drive each defense through hundreds of trials, "
         "produce a comparable security profile so two configurations can be argued about with data."),
    ], font_size=18, bullet_color=ACCENT_RED)


def slide_threat_model(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Threat model", "What we model and what we don't")
    headers = ["Dimension", "In scope", "Out of scope (future work)"]
    rows = [
        ["Attempts", "Sequential HTTP login requests", "Parallel / distributed"],
        ["Source IPs", "Single IP", "Botnets, proxy pools"],
        ["Targets", "One known username", "Credential stuffing, password spraying"],
        ["Wordlist", "Real public corpora (SecLists 10k)", "Adaptive / personalised guesses"],
        ["Knowledge", "Black-box: only HTTP responses", "Insider / source-code visibility"],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(3.6),
              headers, rows, font_size=14)
    add_textbox(s, Inches(0.5), Inches(5.6), w - Inches(1.0), Inches(0.6),
                "Attacker capability tiers - naive bot, PoW-solving bot, human-in-the-loop - "
                "are explicit knobs in the framework.",
                font_size=15, color=TEXT_DARK, font=FONT_BODY)


def slide_section_divider(prs, w, h, number: str, title: str, subtitle: str):
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, w, h, NAVY)
    add_rect(s, 0, h / 2 + Inches(0.1), w, Inches(0.06), ACCENT_RED)
    add_textbox(s, Inches(0.6), h / 2 - Inches(1.4), w - Inches(1.2), Inches(0.8),
                f"Section {number}",
                font_size=22, color=WHITE, font=FONT_HEADING)
    add_textbox(s, Inches(0.6), h / 2 - Inches(0.8), w - Inches(1.2), Inches(1.4),
                title, font_size=44, bold=True, color=WHITE, font=FONT_HEADING)
    add_textbox(s, Inches(0.6), h / 2 + Inches(0.5), w - Inches(1.2), Inches(0.6),
                subtitle, font_size=18, color=NAVY_LIGHT, font=FONT_BODY)


def slide_testbed_arch(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Testbed architecture", "Four pieces, each independently configurable")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(5.0), [
        "Target system: a small Flask login service. Every defense is a tunable knob.",
        "Attacker client: a single password-guessing tool with capability flags for naive bots, "
        "PoW-solving bots, and human-in-the-loop attackers.",
        "Orchestrator: boots a fresh target per configuration, runs the attacker, "
        "aggregates results.",
        "Wordlist source: real public corpora (SecLists), sampled per trial.",
        "Each trial gets an isolated process and a fresh port - no state leaks between "
        "configurations.",
        "Every measurement is seeded and reproducible from the same inputs.",
    ], font_size=17)


def slide_wordlist(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Wordlist methodology", "Distribution, not anecdote")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(2.5), [
        "Old approach: one fixed wordlist with the password at a known position. One data point.",
        "New approach: every trial gets its own randomized wordlist:",
        (1, "Sample N entries from a real public password corpus (SecLists 10k)."),
        (1, "Insert the target password at a uniformly-random position."),
        (1, "Record the seed and target position so the trial is reproducible."),
    ], font_size=17)
    add_rect(s, Inches(0.6), Inches(5.0), w - Inches(1.2), Inches(0.06), ACCENT_GREEN)
    add_bullets(s, Inches(0.6), Inches(5.2), w - Inches(1.2), Inches(1.8), [
        "50 trials per config x 21 configs x 4 wordlist sizes (100 / 200 / 300 / 500) = 3,500+ trial measurements.",
        "Each defense exercised against 50 randomized target depths per wordlist size.",
        "Headline 100-word run finished in 13.7 hours; longer runs were stopped after the layered configs.",
    ], font_size=16, bullet_color=ACCENT_GREEN)


def slide_attacker_tiers(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Attacker capability tiers", "Defenses that beat one tier are bypassed by another")
    headers = ["Tier", "Behaviour"]
    rows = [
        ["Naive scripted bot", "Sends username and password. No JavaScript, no humans, no header spoofing."],
        ["Header-aware bot", "Same as above, but spoofs a normal browser User-Agent."],
        ["PoW-solving bot", "Includes a SHA-256 puzzle solver in its loop."],
        ["Human-in-the-loop", "A human (or commercial solver service) handles CAPTCHAs."],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(3.4),
              headers, rows, font_size=15)
    add_textbox(s, Inches(0.5), Inches(5.4), w - Inches(1.0), Inches(0.8),
                "Same attacker tool models all four tiers - we toggle behaviour per run "
                "to compare each defense across the spectrum.",
                font_size=14, color=TEXT_MUTED, font=FONT_BODY)


def slide_defense_catalog(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Ten defenses, four required categories",
                  "Required project categories highlighted in navy")
    headers = ["Category", "Mechanism", "What it does"]
    rows = [
        [("Account lockout", NAVY), "Account lockout",
         "Freezes the account for a duration after N consecutive failures."],
        [("Rate limiting", NAVY), "IP rate limit",
         "Caps attempts from a single IP in a sliding window."],
        [("Rate limiting", NAVY), "Permanent IP ban",
         "Blacklists the IP after K failures within a long window."],
        [("Progressive delays", NAVY), "Tarpit",
         "Server sleeps a fixed amount before each failed response."],
        [("Progressive delays", NAVY), "IP exponential backoff",
         "Per-IP cooldown that doubles with each failure, with a cap."],
        [("Cost amplification", TEXT_MUTED), "Slow password hash",
         "pbkdf2 or scrypt to inflate per-attempt CPU cost on the server."],
        [("Bot vs human filter", TEXT_MUTED), "Proof-of-work",
         "Server demands a SHA-256 puzzle after N failures."],
        [("Bot vs human filter", TEXT_MUTED), "CAPTCHA challenge",
         "Server demands a human-solvable token after N failures."],
        [("Bot vs human filter", TEXT_MUTED), "Honeypot usernames",
         "Contact with watched usernames (admin, root) triggers an instant ban."],
        [("Bot vs human filter", TEXT_MUTED), "Header anomaly detection",
         "Block requests missing typical browser headers."],
    ]
    add_table(s, Inches(0.4), Inches(1.6), w - Inches(0.8), Inches(5.4),
              headers, rows, font_size=12)


def slide_metrics(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "How they're measured",
                  "Six metrics per config form a comparable security profile")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(5.0), [
        "Breach rate - fraction of trials where the attacker hit the password.",
        "Median time-to-crack - wall-clock seconds, with min/max range.",
        "Effective request rate - requests/sec the attacker could sustain.",
        "Response status mix - 401 / 423 / 429 / 403 by reason "
        "(PoW required, CAPTCHA, account locked, IP banned, anomaly).",
        "First-hit position - where in the wordlist the password was found.",
        "Position vs time scatter - how time-to-crack scales with target depth.",
    ], font_size=18)


def slide_summary_chart(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Summary across 21 configs x 50 trials (100-word run)",
                  "Three regimes: hard block, lottery window, slow-down only")
    img_w = Inches(7.5)
    img_h = Inches(5.0)
    add_image_safe(s, chart("chart_verdict.png"),
                   Inches(0.4), Inches(1.7), img_w, img_h)
    # Side annotations
    nx = Inches(8.2)
    add_textbox(s, nx, Inches(1.9), Inches(4.8), Inches(0.5),
                "Hard block (0%)", font_size=18, bold=True, color=ACCENT_GREEN)
    add_textbox(s, nx, Inches(2.4), Inches(4.8), Inches(0.6),
                "3 of 21: defense fires from request 1 (honeypot, anomaly, exp-backoff).",
                font_size=13, color=TEXT_DARK)
    add_textbox(s, nx, Inches(3.2), Inches(4.8), Inches(0.5),
                "Lottery window (2-4%)", font_size=18, bold=True, color=ACCENT_AMBER)
    add_textbox(s, nx, Inches(3.7), Inches(4.8), Inches(0.7),
                "8 of 21: gated on N=5 fails - breach when password lands in attempts 1-5.",
                font_size=13, color=TEXT_DARK)
    add_textbox(s, nx, Inches(4.7), Inches(4.8), Inches(0.5),
                "Slow-down only (100%)", font_size=18, bold=True, color=ACCENT_RED)
    add_textbox(s, nx, Inches(5.2), Inches(4.8), Inches(0.6),
                "9 of 21: attacker waits, password comes out every trial.",
                font_size=13, color=TEXT_DARK)
    add_textbox(s, nx, Inches(6.0), Inches(4.8), Inches(0.5),
                "Probabilistic (82%)", font_size=18, bold=True, color=ACCENT_AMBER)
    add_textbox(s, nx, Inches(6.5), Inches(4.8), Inches(0.5),
                "F2_pow_22bit - solver sometimes blew its budget.",
                font_size=13, color=TEXT_DARK)


def slide_hard_stoppers(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Hard block (0% breach in 50 trials)",
                  "Defense fires from request 1 - attacker never gets a foothold")
    headers = ["Config", "Median elapsed", "Why it blocks instantly"]
    rows = [
        [("L_honeypot_username", NAVY), "0.5s", "Username 'admin' triggers permanent ban on request 1"],
        [("M_anomaly_no_ua", NAVY), "0.5s", "Missing User-Agent flagged on request 1"],
        [("E_ip_exp_backoff", NAVY), "0.7s", "Per-IP backoff hits cap before any breach window opens"],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(2.0),
              headers, rows, font_size=15)
    add_rect(s, Inches(0.6), Inches(4.4), w - Inches(1.2), Inches(0.06), ACCENT_GREEN)
    add_bullets(s, Inches(0.6), Inches(4.6), w - Inches(1.2), Inches(2.3), [
        "These are the only configs that block before the password could possibly be tried.",
        "Honeypot and anomaly only work because the attacker is naive (admin user / missing UA).",
        "A header-aware attacker bypasses both - see M2_anomaly_normal_ua at 100% breach.",
    ], font_size=15, bullet_color=ACCENT_GREEN)


def slide_lottery_window(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Lottery window (2-4% breach in 50 trials)",
                  "Defenses gated on N=5 failures: breach when password lands in attempts 1-5")
    headers = ["Config", "Breach", "Median elapsed", "Defense threshold"]
    rows = [
        [("B_account_lockout", NAVY), ("2%", ACCENT_AMBER), "1.3s", "5 fails -> 60s lock"],
        [("G_pow_naive_attacker", NAVY), ("2%", ACCENT_AMBER), "1.3s", "5 fails -> PoW (no solver)"],
        [("J_captcha_naive", NAVY), ("2%", ACCENT_AMBER), "1.3s", "5 fails -> CAPTCHA (no solver)"],
        [("I_perma_ban", NAVY), ("2%", ACCENT_AMBER), "1.8s", "8 fails -> permanent ban"],
        [("C_ip_rate_limit", NAVY), ("4%", ACCENT_AMBER), "2.1s", "10 attempts / 30s window"],
        [("H_layered_basic", NAVY), ("2%", ACCENT_AMBER), "26.8s", "lockout dominates the stack"],
        [("H3_full_stack", NAVY), ("4%", ACCENT_AMBER), "30.1s", "every mechanism enabled"],
    ]
    add_table(s, Inches(0.4), Inches(1.7), w - Inches(0.8), Inches(3.4),
              headers, rows, font_size=13)
    add_rect(s, Inches(0.6), Inches(5.4), w - Inches(1.2), Inches(0.06), ACCENT_AMBER)
    add_bullets(s, Inches(0.6), Inches(5.6), w - Inches(1.2), Inches(1.8), [
        "P(breach) ~= 5/N where N is the wordlist size - the gating window is a structural blind spot.",
        "Median first-success position across all lottery breaches: 3.",
        "Pair with a defense that's active from request 1 (slow hash, honeypot) to cover the window.",
    ], font_size=14, bullet_color=ACCENT_AMBER)


def slide_slow_downs(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Slow-down only (100% breach in 50 trials)",
                  "Cost the attacker time, but the password came out every trial")
    headers = ["Config", "Median time", "Slowdown vs baseline"]
    rows = [
        [("A_baseline", TEXT_DARK), "14.5s", "1.0x"],
        [("K2_slow_hash_scrypt", ACCENT_RED), "14.4s", ("1.0x - no slowdown", ACCENT_RED)],
        [("M2_anomaly_normal_ua", TEXT_DARK), "14.4s", "1.0x - bypassed by spoofed UA"],
        [("J2_captcha_human", TEXT_DARK), "14.5s", "1.0x"],
        [("F_pow_smart_attacker (18-bit)", TEXT_DARK), "34.4s", "2.4x"],
        [("D_tarpit_500ms", TEXT_DARK), "64.0s", "4.4x"],
        [("K_slow_hash_pbkdf2 (600k iters)", TEXT_DARK), "83.5s", "5.8x"],
        [("D2_tarpit_1s", TEXT_DARK), "113.5s", "7.8x"],
        [("D3_tarpit_2s", TEXT_DARK), "212.5s", ("14.7x", ACCENT_GREEN)],
        [("F2_pow_22bit (smart, 22-bit)", ACCENT_AMBER), "305.9s (82% breach)", "21x"],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(4.8),
              headers, rows, font_size=13)
    add_textbox(s, Inches(0.5), Inches(6.6), w - Inches(1.0), Inches(0.6),
                "Tarpit and slow-hash slowdowns are linear in wordlist depth - predictable and tunable. (Next slide.)",
                font_size=14, color=TEXT_MUTED, font=FONT_BODY)


def slide_finding_scrypt(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Finding 1: scrypt was barely a defense",
                  "Default werkzeug scrypt:32768:8:1 was statistically indistinguishable from baseline")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(2.5), [
        "scrypt:32768:8:1 - 14.40s (vs 14.46s baseline). Lost in HTTP roundtrip noise.",
        "Server-side scrypt cost dominated by n=32768 (~16ms on test hardware).",
        "pbkdf2:sha256:600000 was meaningfully slower (83.5s, 5.8x baseline).",
    ], font_size=18)
    add_rect(s, Inches(0.6), Inches(4.5), w - Inches(1.2), Inches(0.06), ACCENT_RED)
    add_textbox(s, Inches(0.6), Inches(4.7), w - Inches(1.2), Inches(2.0),
                "Lesson: \"we use a slow hash\" is not a defense unless parameters are tuned for the "
                "target hardware. Production deployments should benchmark and aim for ~100 ms per verify. "
                "Don't trust library defaults.",
                font_size=18, color=ACCENT_RED, bold=True)


def slide_finding_pow(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Finding 2: PoW sits on a probabilistic cliff",
                  "Past a difficulty threshold, PoW shifts from slow-down to probabilistic block")
    headers = ["Config", "Difficulty", "Breach rate (50 trials)", "Median time"]
    rows = [
        [("F_pow_smart_attacker", TEXT_DARK), "18-bit", "100%", "34.4s"],
        [("F2_pow_22bit", ACCENT_AMBER), "22-bit", ("82%", ACCENT_AMBER), "305.9s"],
    ]
    add_table(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(1.8),
              headers, rows, font_size=15)
    add_bullets(s, Inches(0.6), Inches(4.0), w - Inches(1.2), Inches(2.8), [
        "At 22-bit difficulty the attacker's 5M-attempt-per-puzzle budget runs out on roughly "
        "1 in 5 trials.",
        "Defense moves from slow-down to probabilistic block.",
        "At 200-word the same config breached 30/33 (91%) - more depth means more puzzles, but "
        "the solver kept squeaking through.",
        "Cliff is a function of solver-budget per puzzle, not just difficulty bits.",
    ], font_size=16)


def slide_finding_scaling(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Finding 3: time scales linearly with wordlist depth",
                  "Each line's slope is the per-attempt cost the defense imposes")
    add_image_safe(s, chart("chart_wordlist_scaling.png"),
                   Inches(0.4), Inches(1.6), Inches(8.6), Inches(5.0))
    nx = Inches(9.2)
    add_textbox(s, nx, Inches(1.8), Inches(4.0), Inches(0.5),
                "Per-attempt cost", font_size=18, bold=True, color=NAVY)
    add_bullets(s, nx, Inches(2.4), Inches(4.0), Inches(4.5), [
        "D3_tarpit_2s: ~2.0 s / attempt",
        "D2_tarpit_1s: ~1.0 s / attempt",
        "K_slow_hash_pbkdf2: ~0.83 s / attempt",
        "D_tarpit_500ms: ~0.5 s / attempt",
        "F_pow_smart_attacker: ~0.4 s / attempt",
        "A_baseline / K2_scrypt: overlap - scrypt added zero observable cost",
    ], font_size=13)
    add_textbox(s, nx, Inches(6.0), Inches(4.0), Inches(1.0),
                "Buying time per attempt: budget = cost x depth.",
                font_size=14, color=TEXT_MUTED, font=FONT_BODY)


def slide_finding_breach_scaling(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Finding 4: lottery breaches track 5/N",
                  "Defenses gated on N=5 failures leak when password lands in attempts 1-5")
    add_image_safe(s, chart("chart_breach_scaling.png"),
                   Inches(0.4), Inches(1.6), Inches(8.6), Inches(4.6))
    nx = Inches(9.2)
    add_bullets(s, nx, Inches(1.8), Inches(4.0), Inches(5.0), [
        "100-word: 5/100 = 5% expected, 2-4% observed",
        "200-word: 5/200 = 2.5%, 2% observed",
        "300-word: 5/300 = 1.7%, 0% observed (sampling noise)",
        "500-word: 5/500 = 1%, 2-4% observed for some",
        "I_perma_ban gated on 8 fails - 8/N decays slower",
    ], font_size=13)
    add_rect(s, Inches(0.6), Inches(6.5), w - Inches(1.2), Inches(0.06), ACCENT_AMBER)
    add_textbox(s, Inches(0.6), Inches(6.7), w - Inches(1.2), Inches(0.6),
                "Operational fix: pair gated defenses with one that's active from request 1 "
                "(slow hash, honeypot, anomaly).",
                font_size=14, color=TEXT_DARK, font=FONT_BODY, bold=True)


def slide_layered(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Layered defense converges",
                  "Account lockout dominates - the rest are insurance for when it fails")
    headers = ["Config", "Defenses", "Median time", "Breach"]
    rows = [
        ["H_layered_basic", "lockout + rate-limit + tarpit + PoW", "26.8s", ("2%", ACCENT_AMBER)],
        ["H2_layered_with_ban", "+ perma-ban + slow hash", "30.4s", ("2%", ACCENT_AMBER)],
        ["H3_full_stack", "+ CAPTCHA + honeypot + anomaly", "30.1s", ("4%", ACCENT_AMBER)],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(2.4),
              headers, rows, font_size=14)
    add_rect(s, Inches(0.6), Inches(4.5), w - Inches(1.2), Inches(0.06), NAVY)
    add_bullets(s, Inches(0.6), Inches(4.7), w - Inches(1.2), Inches(2.5), [
        "Wall-clock is dominated by account lockout - other defenses never get exercised.",
        "The 2-4% breach is the same lottery-window effect: password lands in attempts 1-5 before "
        "lockout activates.",
        ("That's a feature, not a bug: the cheapest defense wins, the rest are insurance for "
         "when it fails or is misconfigured."),
        "Defense-in-depth is about failure modes, not steady-state performance.",
    ], font_size=15)


def slide_recommendations(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "What to actually deploy",
                  "Minimum-viable stack - three layers with tuned parameters")
    rows = [
        [("1. Account lockout", ACCENT_GREEN),
         "~5 failures -> ≥60s lockout. Reset on legitimate login. Cheap, hard-stops named-target attacks."],
        [("2. IP-based throttling", ACCENT_GREEN),
         "Sliding window or exponential backoff. Cap, don't permanent-ban (avoid IP-reuse pain)."],
        [("3. Slow password hash", ACCENT_GREEN),
         "argon2id or pbkdf2, tuned to ~100ms/verify on prod hardware. Verify cost - defaults are weak."],
    ]
    add_table(s, Inches(0.5), Inches(1.7), w - Inches(1.0), Inches(2.6),
              ["Layer", "Configuration"], rows, font_size=14)

    add_rect(s, Inches(0.6), Inches(4.6), w - Inches(1.2), Inches(0.06), ACCENT_RED)
    add_bullets(s, Inches(0.6), Inches(4.8), w - Inches(1.2), Inches(2.0), [
        "Don't rely on (against motivated attackers): tarpits alone, default-parameter scrypt, "
        "header anomaly checks (trivially defeated by a User-Agent string), low-bit PoW.",
        "Useful add-ons: honeypot usernames, geographic anomaly scoring, MFA on sensitive accounts.",
    ], font_size=15, bullet_color=ACCENT_RED)


def slide_reproduce(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Reproducing this study",
                  "Every artifact lands under login-lab/logs/benchmark/<UTC stamp>/")
    add_code_block(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(3.5), """\
git clone <repo>
cd Research-Project
python -m venv .venv && .venv/Scripts/activate
pip install -r requirements.txt

# Refresh wordlists (uses existing download scripts)
python passwords/download-scripts/run_all.py

# Full statistical run (~60 min on a laptop)
python scripts/benchmark_defenses.py --trials 5 --seed 1337

# Charts & HTML report
python scripts/make_charts.py --open""", font_size=14)
    add_textbox(s, Inches(0.6), Inches(5.6), w - Inches(1.2), Inches(0.5),
                "Outputs: per-attempt CSVs, server logs, generated wordlists, chart PNGs, HTML/MD report.",
                font_size=14, color=TEXT_MUTED, font=FONT_BODY)


def slide_future(prs, w, h):
    s = add_blank_slide(prs)
    add_title_bar(s, w, "Future work", "Where the framework can be extended")
    add_bullets(s, Inches(0.6), Inches(1.7), w - Inches(1.2), Inches(5.0), [
        "Distributed attacker - multi-IP, multi-process to expose IP-only defenses.",
        "Cross-system comparison - point the same attack client at "
        "WordPress, Authelia, Gitea, Keycloak. The framework already speaks plain HTTP.",
        "Adaptive attackers - observe response timing and codes, switch strategy mid-run.",
        "Real CAPTCHA / MFA endpoints - replace magic-token stubs with hCaptcha or TOTP.",
        "Cost modeling - translate \"13x slowdown\" into dollar cost per credential at "
        "cloud-attacker rates.",
    ], font_size=18)


def slide_questions(prs, w, h):
    s = add_blank_slide(prs)
    add_rect(s, 0, 0, w, h, NAVY)
    add_textbox(s, Inches(0.6), h / 2 - Inches(0.9), w - Inches(1.2), Inches(2.0),
                "Questions?", font_size=96, bold=True, color=WHITE,
                font=FONT_HEADING, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── main ─────────────────────────────────────────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    w, h = prs.slide_width, prs.slide_height

    # Build all slides
    slides_to_build = [
        ("title", lambda: slide_title(prs, w, h)),
        ("problem", lambda: slide_problem(prs, w, h)),
        ("threat", lambda: slide_threat_model(prs, w, h)),
        ("sec1", lambda: slide_section_divider(prs, w, h, "1", "Testbed",
                                               "How we measure each defense reproducibly")),
        ("arch", lambda: slide_testbed_arch(prs, w, h)),
        ("wordlist", lambda: slide_wordlist(prs, w, h)),
        ("attacker", lambda: slide_attacker_tiers(prs, w, h)),
        ("sec2", lambda: slide_section_divider(prs, w, h, "2", "Protections",
                                               "Ten mechanisms across four categories")),
        ("catalog", lambda: slide_defense_catalog(prs, w, h)),
        ("metrics", lambda: slide_metrics(prs, w, h)),
        ("sec3", lambda: slide_section_divider(prs, w, h, "3", "Results",
                                               "What blocked, what just slowed the attacker down")),
        ("summary", lambda: slide_summary_chart(prs, w, h)),
        ("hard", lambda: slide_hard_stoppers(prs, w, h)),
        ("lottery", lambda: slide_lottery_window(prs, w, h)),
        ("slow", lambda: slide_slow_downs(prs, w, h)),
        ("scrypt", lambda: slide_finding_scrypt(prs, w, h)),
        ("pow", lambda: slide_finding_pow(prs, w, h)),
        ("scaling", lambda: slide_finding_scaling(prs, w, h)),
        ("breach_scaling", lambda: slide_finding_breach_scaling(prs, w, h)),
        ("layered", lambda: slide_layered(prs, w, h)),
        ("sec4", lambda: slide_section_divider(prs, w, h, "4", "Recommendations",
                                               "What to actually deploy")),
        ("recs", lambda: slide_recommendations(prs, w, h)),
        ("future", lambda: slide_future(prs, w, h)),
        ("questions", lambda: slide_questions(prs, w, h)),
    ]

    for _name, builder in slides_to_build:
        builder()

    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")
    print(f"Slide count: {len(prs.slides)}")
    if LATEST is None:
        print("(no benchmark run dir found - chart slides will show placeholders)")
    else:
        print(f"Charts pulled from: {LATEST}")


if __name__ == "__main__":
    main()
